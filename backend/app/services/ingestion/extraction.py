import csv
import json
import zipfile
from io import BytesIO, StringIO

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def extract_text(file_type: str, data: bytes, ocr_text: str | None = None) -> tuple[str, int | None]:
    if ocr_text:
        return ocr_text, None

    if file_type == "pdf":
        reader = PdfReader(BytesIO(data))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(pages), len(reader.pages)

    if file_type == "docx":
        doc = DocxDocument(BytesIO(data))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip()), None

    if file_type in {"xlsx", "xls"}:
        wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                rows.append("\t".join("" if cell is None else str(cell) for cell in row))
        return "\n".join(rows), None

    if file_type == "pptx":
        pres = Presentation(BytesIO(data))
        texts = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n\n".join(texts), len(pres.slides)

    if file_type in {"txt", "md", "csv", "json", "xml"}:
        text = data.decode("utf-8", errors="ignore")
        if file_type == "csv":
            reader = csv.reader(StringIO(text))
            text = "\n".join(", ".join(row) for row in reader)
        return text, None

    if file_type in {"html", "htm"}:
        soup = BeautifulSoup(data, "lxml")
        return soup.get_text("\n"), None

    if file_type in {"png", "jpg", "jpeg", "tiff", "tif"}:
        return "", None

    if file_type == "zip":
        texts = []
        with zipfile.ZipFile(BytesIO(data)) as archive:
            for name in archive.namelist():
                if name.endswith((".txt", ".md", ".csv", ".json", ".xml", ".html")):
                    texts.append(archive.read(name).decode("utf-8", errors="ignore"))
        return "\n\n".join(texts), None

    return data.decode("utf-8", errors="ignore"), None
